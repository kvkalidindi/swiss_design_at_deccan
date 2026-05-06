"""Build deccan.potx using python-pptx.

PowerPoint template with the four pieces of brand "furniture":
  1. Cover slide  — large logo + title + metadata block
  2. (no top header on slides — cover slide carries that role)
  3. Slide footer — small logo + section name + slide number on every body slide
  4. End slide    — centered logo + brand line

See skill/references/document-furniture.md for the canonical spec.
"""
from __future__ import annotations

import json
import zipfile
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from scripts.lib.office_theme import build_theme_xml

ROOT = Path(__file__).resolve().parents[2]
PALETTE = ROOT / "outputs" / "palette.json"
LOGO = ROOT / "data" / "logo.png"
OUT = ROOT / "office" / "templates" / "deccan.potx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _palette() -> dict:
    return json.loads(PALETTE.read_text(encoding="utf-8"))


def _add_text(slide, text: str, left, top, width, height, *,
              size: int, color_hex: str, bold: bool = False,
              align=PP_ALIGN.LEFT, font: str = "IBM Plex Sans") -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color_hex)


def _add_slide_footer(slide, blue_500: str, section_name: str, slide_number_text: str) -> None:
    """Add the persistent slide footer: logo (left) + section (center) + page # (right) +
    a thin accent rule above."""
    # Bottom rule (thin line) at 30% blue-500 — render as a 1pt rectangle.
    rule_top = SLIDE_H - Inches(0.55)
    rule_left = Inches(0.5)
    rule_right_margin = Inches(0.5)
    rule_width = SLIDE_W - rule_left - rule_right_margin
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rule_left, rule_top, rule_width, Emu(9525))
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor.from_string(blue_500)
    rule.fill.fore_color.brightness = 0.7  # lighten to ~30%-ish presence

    if LOGO.exists():
        slide.shapes.add_picture(
            str(LOGO),
            Inches(0.5),
            SLIDE_H - Inches(0.45),
            width=Inches(0.9),
        )

    _add_text(
        slide, section_name,
        SLIDE_W / 2 - Inches(2),
        SLIDE_H - Inches(0.4),
        Inches(4),
        Inches(0.3),
        size=9, color_hex="595959", align=PP_ALIGN.CENTER,
    )

    _add_text(
        slide, slide_number_text,
        SLIDE_W - Inches(2.0),
        SLIDE_H - Inches(0.4),
        Inches(1.5),
        Inches(0.3),
        size=9, color_hex="595959", align=PP_ALIGN.RIGHT,
    )


def _build_cover_slide(prs, blue_500: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if LOGO.exists():
        slide.shapes.add_picture(
            str(LOGO), Inches(0.7), Inches(0.7), width=Inches(2.5),
        )

    _add_text(
        slide, "Document Title",
        Inches(0.7), Inches(2.6),
        Inches(11.5), Inches(1.2),
        size=44, color_hex=blue_500, bold=False, align=PP_ALIGN.LEFT,
    )

    _add_text(
        slide, "Subtitle / one-line summary",
        Inches(0.7), Inches(3.7),
        Inches(11.5), Inches(0.6),
        size=20, color_hex="595959", align=PP_ALIGN.LEFT,
    )

    # Accent rule under subtitle
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(4.4),
        Inches(11.5), Emu(12700),
    )
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor.from_string(blue_500)

    # Metadata block (caption-style key/value pairs, left-aligned bottom)
    labels = ["DOCUMENT TYPE", "PREPARED BY", "DATE", "VERSION", "CLASSIFICATION"]
    values = ["Presentation", "Author Name", "YYYY-MM-DD", "v1.0", "Confidential"]
    y = Inches(4.7)
    for label, value in zip(labels, values):
        _add_text(
            slide, label, Inches(0.7), y, Inches(2.2), Inches(0.25),
            size=8, color_hex="8A8786", bold=False, align=PP_ALIGN.LEFT,
        )
        _add_text(
            slide, value, Inches(2.9), y, Inches(8.0), Inches(0.25),
            size=10, color_hex="1C1917", align=PP_ALIGN.LEFT,
        )
        y += Inches(0.34)


def _build_section_divider(prs, blue_500: str, section_name: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _add_text(
        slide, section_name,
        Inches(0.7), Inches(3.0),
        Inches(11.5), Inches(1.2),
        size=44, color_hex=blue_500, bold=False, align=PP_ALIGN.LEFT,
    )
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(4.2),
        Inches(2.0), Emu(12700),
    )
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor.from_string(blue_500)

    _add_slide_footer(slide, blue_500, section_name, "")


def _build_content_slide(prs, blue_500: str, title_text: str, slide_idx: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _add_text(
        slide, title_text,
        Inches(0.5), Inches(0.4),
        Inches(12.3), Inches(0.7),
        size=28, color_hex=blue_500, bold=False, align=PP_ALIGN.LEFT,
    )

    _add_text(
        slide, "Body content area. Replace this with your slide content.",
        Inches(0.5), Inches(1.4),
        Inches(12.3), Inches(5.0),
        size=18, color_hex="1C1917", align=PP_ALIGN.LEFT,
    )

    _add_slide_footer(slide, blue_500, "Section name", f"{slide_idx} / {total}")


def _build_end_slide(prs, blue_500: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if LOGO.exists():
        # Centered horizontally, slightly above vertical center
        logo_w = Inches(2.5)
        slide.shapes.add_picture(
            str(LOGO),
            (SLIDE_W - logo_w) / 2,
            Inches(2.6),
            width=logo_w,
        )

    _add_text(
        slide, "Deccan Chemicals",
        Inches(0), Inches(4.6),
        SLIDE_W, Inches(0.5),
        size=20, color_hex="1C1917", align=PP_ALIGN.CENTER,
    )

    _add_text(
        slide, "deccanchemicals.com · Hyderabad, India",
        Inches(0), Inches(5.2),
        SLIDE_W, Inches(0.4),
        size=12, color_hex="8A8786", align=PP_ALIGN.CENTER,
    )


def emit_potx() -> Path:
    palette = _palette()
    blue_500 = palette["blue"]["500"]["hex"].lstrip("#")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover slide
    _build_cover_slide(prs, blue_500)

    # 2. Section divider (example)
    _build_section_divider(prs, blue_500, "Section name")

    # 3-5. Content slide examples (with persistent slide footer)
    _build_content_slide(prs, blue_500, "Content slide title", 3, 5)
    _build_content_slide(prs, blue_500, "Two-column slide title", 4, 5)
    _build_content_slide(prs, blue_500, "Title-only slide", 5, 5)

    # 6. End slide
    _build_end_slide(prs, blue_500)

    tmp = OUT.with_suffix(".pptx")
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(tmp)

    _replace_theme_in_pptx(tmp, build_theme_xml(palette))
    _convert_to_potx(tmp, OUT)
    tmp.unlink()
    return OUT


def _replace_theme_in_pptx(path: Path, theme_xml: str) -> None:
    """Open the .pptx zip and replace ppt/theme/theme1.xml with our theme."""
    buf = BytesIO()
    with zipfile.ZipFile(path, "r") as zin:
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item == "ppt/theme/theme1.xml":
                    data = theme_xml.encode("utf-8")
                zout.writestr(item, data)
    path.write_bytes(buf.getvalue())


def _convert_to_potx(src: Path, dst: Path) -> None:
    """Convert .pptx package to .potx by updating [Content_Types].xml."""
    buf = BytesIO()
    with zipfile.ZipFile(src, "r") as zin:
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item == "[Content_Types].xml":
                    text = data.decode("utf-8")
                    text = text.replace(
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                    )
                    data = text.encode("utf-8")
                zout.writestr(item, data)
    dst.write_bytes(buf.getvalue())
